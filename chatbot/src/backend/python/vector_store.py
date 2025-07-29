import os
import re
import mimetypes
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

import ftfy
import fitz
import spacy
import win32com.client
from PIL import Image, ImageGrab
import imagehash
from pptx import Presentation
from docx import Document
from docx.oxml.ns import qn

from dotenv import load_dotenv

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import TextLoader, PyPDFLoader
from langchain_huggingface import HuggingFaceEmbeddings

from langchain.schema import Document as langDocument
from langchain_groq import ChatGroq
from langchain_core.output_parsers import StrOutputParser

import os
import re
import ftfy
from pathlib import Path
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings

from langchain.schema import Document as langDocument
from langchain_groq import ChatGroq
from langchain_core.output_parsers import StrOutputParser

## define root directory and subsequent sub directories
root_dir = Path(os.getcwd())  # Get the current working directory (the root of your repository)
data_dir = root_dir / "data"  # Gets the "data" folder relative to the root directory
images_dir = data_dir / "images"
pdf_dir = data_dir / "pdf"
docx_dir = data_dir / "word"
pptx_dir = data_dir / "pptx"
cleaned_dir=data_dir/ "cleaned"
vertztec_collection= data_dir/ "verztec_logo"

# Load environment variables and spaCy model
load_dotenv()
nlp = spacy.load("en_core_web_sm")

# Initialize models
api_key = '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'  # Set in your .env file
model = "deepseek-r1-distill-llama-70b"
deepseek = ChatGroq(api_key=api_key, model_name=model)
deepseek_chain = deepseek | StrOutputParser()

embedding_model = HuggingFaceEmbeddings(
    model_name="BAAI/bge-large-en-v1.5",
    encode_kwargs={'normalize_embeddings': True}
)   

# Text splitters
fallback_splitter = RecursiveCharacterTextSplitter(chunk_size=250, chunk_overlap=50)

def spacy_fallback_splitter(text, chunk_size=400, overlap=50):
    doc = nlp(text)
    sentences = [sent.text.strip() for sent in doc.sents if sent.text.strip()]
    chunks, current_chunk = [], ""
    for sentence in sentences:
        if len(current_chunk) + len(sentence) <= chunk_size:
            current_chunk += " " + sentence
        else:
            chunks.append(current_chunk.strip())
            current_chunk = sentence
    if current_chunk:
        chunks.append(current_chunk.strip())

    if overlap > 0:
        overlapped_chunks = []
        for i in range(len(chunks)):
            start = max(0, i - 1)
            merged = " ".join(chunks[start:i+1])
            overlapped_chunks.append(merged.strip())
        return overlapped_chunks

    return chunks

def replace_image_tags_with_placeholders(text):
    pattern = r"<\|image_start\|>(.*?)<\|image_end\|>"
    matches = re.findall(pattern, text)
    replacements = {}
    image_name_map = {}
    for i, match in enumerate(matches):
        key = f"[[IMAGE_TAG_{i}]]"
        full_tag = f"<|image_start|>{match}<|image_end|>"
        text = text.replace(full_tag, key)
        replacements[key] = full_tag
        image_name_map[key] = match  # Just the image file name
    return text, replacements, image_name_map

def restore_image_placeholders_and_collect_metadata(chunks, replacements, image_name_map):
    restored_chunks = []
    chunk_image_metadata = []

    for chunk in chunks:
        images_in_chunk = []
        for placeholder, original in replacements.items():
            if placeholder in chunk:
                chunk = chunk.replace(placeholder, original)
                images_in_chunk.append(image_name_map[placeholder])
        restored_chunks.append(chunk)
        chunk_image_metadata.append(images_in_chunk)
    
    return restored_chunks, chunk_image_metadata

def get_faiss_index_path(country, department, base_dir="faiss_indices"):
    """
    Generate FAISS index path based on country and department.
    Creates directory structure: faiss_indices/{country}/{department}/
    Special case: admin/master -> faiss_indices/admin_master/
    """
    # Normalize inputs
    country = country.lower().strip()
    department = department.lower().strip()
    
    # Special case for admin master index
    if country == 'admin' and department == 'master':
        index_dir = Path(base_dir) / "admin_master"
    else:
        # Create directory path
        index_dir = Path(base_dir) / country / department
    
    index_dir.mkdir(parents=True, exist_ok=True)
    
    return str(index_dir / "faiss_index")

def load_single_file(file_path, embedding_model, faiss_index_path=None, country=None, department=None):
    """
    Processes a single file, generates embeddings, splits text into chunks,
    and updates the FAISS index if given, or creates a new one if no index is provided.
    Now supports country and department-specific FAISS indices.
    """

    # Ensure the file is valid
    file_path = Path(file_path)
    if not file_path.exists():
        print(f"[ERROR] File not found: {file_path}")
        return None

    # If country and department are provided, generate the appropriate FAISS path
    if country and department and not faiss_index_path:
        faiss_index_path = get_faiss_index_path(country, department)
        print(f"[INFO] Using country/department-specific FAISS index: {faiss_index_path}")

    base_name = file_path.stem
    cleaned_filename = base_name.lower().replace(" ", "_") + ".txt"  # Cleaned filename

    # Read the file content
    with open(file_path, 'r', encoding='utf-8') as f:
        text = f.read().lower()

    # Generate semantic description with LLM (Deepseek)
    query = (
        "Summarize this internal document in 50 words. Describe its purpose, main contents, "
        "and types of internal questions it can help answer. "
        "Respond concisely and in a single paragraph:\n\n" + text
    )
    raw_response = deepseek_chain.invoke(query)
    description = re.sub(r"<think>.*?</think>", "", raw_response, flags=re.DOTALL).strip()

    # Replace image tags with placeholders
    text, image_replacements, image_name_map = replace_image_tags_with_placeholders(text)

    # Split text into chunks
    smart_chunks = spacy_fallback_splitter(text)

    # Restore image tags and get images per chunk
    smart_chunks, chunk_image_lists = restore_image_placeholders_and_collect_metadata(
        smart_chunks, image_replacements, image_name_map
    )

    # Build langDocument objects for new chunks
    all_chunks = []
    for i, chunk in enumerate(smart_chunks):
        image_list = chunk_image_lists[i]

        # Add country and department metadata if provided
        chunk_metadata = {
            "source": base_name,
            "chunk_id": f"{base_name}_{i}",
            "clean_chunk": chunk,
            "images": image_list
        }
        
        if country:
            chunk_metadata["country"] = country.lower()
        if department:
            chunk_metadata["department"] = department.lower()

        enriched_chunk = f"[Description: {description}] [Document: {base_name}] {chunk}"
        all_chunks.append(langDocument(
            page_content=enriched_chunk,
            metadata=chunk_metadata
        ))

    # Check if FAISS index exists
    if faiss_index_path and os.path.exists(faiss_index_path):
        print(f"[INFO] Loading existing FAISS index from {faiss_index_path}")
        faiss_db = FAISS.load_local(
            faiss_index_path, 
            embeddings=embedding_model,
            allow_dangerous_deserialization=True  # Add this parameter
        )
        # Add the new chunks to the existing FAISS index
        faiss_db.add_documents(all_chunks)
    else:
        # Create a new FAISS index
        print(f"[INFO] Creating new FAISS index at {faiss_index_path or f'faiss_index_{base_name}'}")
        faiss_db = FAISS.from_documents(all_chunks, embedding_model)

    # Save updated FAISS index
    save_path = faiss_index_path or f"faiss_index_{base_name}"
    faiss_db.save_local(save_path)

    print(f"[INFO] Processed and updated index for: {cleaned_filename}")
    if country and department:
        print(f"[INFO] Document categorized under: {country.upper()}/{department.upper()}")

    return faiss_db

#-------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
from Documents_Totext import process_single_file

def unified_document_pipeline(file_path, images_dir, cleaned_dir, vertztec_collection, faiss_index_path, embedding_model, country=None, department=None):
    """
    Unified pipeline to process documents and update FAISS index:
    1. Extract and clean document content
    2. Process images and text
    3. Generate embeddings and update FAISS index
    
    Args:
        file_path: Path to input document
        images_dir: Directory for extracted images
        cleaned_dir: Directory for cleaned text
        vertztec_collection: Directory with Verztec logos
        faiss_index_path: Path to FAISS index (optional if country/department provided)
        embedding_model: HuggingFace embedding model
        country: Country for document categorization (e.g., 'china', 'singapore')
        department: Department for document categorization (e.g., 'hr', 'it')
    
    Returns:
        dict: Processing results including FAISS update status
    """
    try:
        # Step 1: Process document and extract content
        processing_result = process_single_file(
            file_path=file_path,
            images_dir=images_dir,
            cleaned_dir=cleaned_dir,
            vertztec_collection=vertztec_collection
        )
        
        if not processing_result["success"]:
            print(f"[ERROR] Document processing failed: {processing_result['error']}")
            return processing_result

        # Step 2: Determine FAISS index path based on country/department
        if country and department and not faiss_index_path:
            faiss_index_path = get_faiss_index_path(country, department)
        elif country and department:
            # If both are provided, use the country/department structure anyway
            faiss_index_path = get_faiss_index_path(country, department)

        # Step 3: Update FAISS index with processed content
        cleaned_text_path = processing_result["cleaned_text_path"]
        faiss_db = load_single_file(
            file_path=cleaned_text_path,
            embedding_model=embedding_model,
            faiss_index_path=faiss_index_path,
            country=country,
            department=department
        )
        
        if faiss_db is None:
            processing_result["success"] = False
            processing_result["error"] = "Failed to update FAISS index"
            return processing_result

        # Return processing details for the Node.js backend to handle database operations
        return {
            **processing_result,
            "faiss_index_path": str(faiss_index_path),
            "country": country,
            "department": department,
            "chunks": [doc.page_content for doc in faiss_db.docstore._dict.values()],
            "metadata": [doc.metadata for doc in faiss_db.docstore._dict.values()]
        }

    except Exception as e:
        error_msg = f"Pipeline failed: {str(e)}"
        print(f"[ERROR] {error_msg}")
        return {
            "original_path": str(file_path),
            "cleaned_text_path": None,
            "success": False,
            "error": error_msg,
            "country": country,
            "department": department
        }

def list_available_indices():
    """
    List all available FAISS indices organized by country and department.
    Returns a dictionary structure showing available combinations.
    """
    base_dir = Path("faiss_indices")
    if not base_dir.exists():
        return {}
    
    indices = {}
    for country_dir in base_dir.iterdir():
        if country_dir.is_dir():
            country = country_dir.name
            indices[country] = []
            
            for dept_dir in country_dir.iterdir():
                if dept_dir.is_dir() and (dept_dir / "faiss_index.index").exists():
                    indices[country].append(dept_dir.name)
    
    return indices

def get_faiss_db_for_query(country=None, department=None, embedding_model=embedding_model):
    """
    Load the appropriate FAISS database for querying based on country and department.
    If no specific country/department is provided, loads the master index.
    """
    if country and department:
        faiss_path = get_faiss_index_path(country, department)
        if os.path.exists(faiss_path):
            try:
                return FAISS.load_local(
                    faiss_path, 
                    embeddings=embedding_model,
                    allow_dangerous_deserialization=True
                )
            except Exception as e:
                print(f"[WARNING] Failed to load FAISS index for {country}/{department}: {e}")
                return None
    
    # Fallback to master index
    current_dir = Path(__file__).parent
    faiss_index_path = current_dir / "faiss_master_index2"
    if os.path.exists(faiss_index_path):
        try:
            return FAISS.load_local(
                str(faiss_index_path), 
                embeddings=embedding_model,
                allow_dangerous_deserialization=True
            )
        except Exception as e:
            print(f"[ERROR] Failed to load master FAISS index: {e}")
            return None
    
    return None

# Example usage:
if __name__ == "__main__":
    # File to process
    file_to_process = r"C:\Users\ethan\OneDrive\Desktop\ICP_Verztec_Chatbot-6\data\pdf\example.pdf"
    
    # Example: Process document for Singapore HR department
    result = unified_document_pipeline(
        file_path=file_to_process,
        images_dir=images_dir,
        cleaned_dir=cleaned_dir,
        vertztec_collection=vertztec_collection,
        faiss_index_path=None,  # Will be auto-generated based on country/department
        embedding_model=embedding_model,
        country="singapore",
        department="hr"
    )
    
    if result["success"]:
        print(f"Pipeline completed successfully!")
        print(f"Document added to: {result['country']}/{result['department']} index")
    else:
        print(f"Pipeline failed: {result['error']}")
    
    # List all available indices
    print("\nAvailable FAISS indices:")
    indices = list_available_indices()
    for country, departments in indices.items():
        print(f"  {country.upper()}: {', '.join(departments)}")