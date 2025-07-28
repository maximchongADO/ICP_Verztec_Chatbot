import pymupdf
import re
import ftfy
import os 
from pathlib import Path
from docx import Document
from PIL import ImageGrab
import win32com.client
import mimetypes
from pptx import Presentation
import os
from pathlib import Path
import mimetypes
import fitz  # pymupdf
import os
from docx import Document
from docx.oxml.ns import qn
from docx.opc.constants import RELATIONSHIP_TYPE as RT
import zipfile

import zipfile
import xml.etree.ElementTree as ET
import os
import docx
from docx import Document
from PIL import Image
import imagehash
import os


## define root directory and subsequent sub directories
root_dir = Path(os.getcwd())  # Get the current working directory (the root of your repository)
data_dir = root_dir / "data"  # Gets the "data" folder relative to the root directory
images_dir = data_dir / "images"
pdf_dir = data_dir / "pdf"
docx_dir = data_dir / "word"
pptx_dir = data_dir / "pptx"
cleaned_dir=data_dir/ "cleaned"
vertztec_collection= data_dir/ "verztec_logo"


def get_all_file_paths(directory, extensions=[".pdf", ".doc", ".docx",".txt",".pptx"]):
    file_paths = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            # Filter files by the given extensions (e.g., PDF, DOC, DOCX files)
            if any(file.lower().endswith(ext) for ext in extensions):
                absolute_path = os.path.abspath(os.path.join(root, file))
                file_paths.append(absolute_path)
    return file_paths



def clean_filename_for_embedding(file_path):
    filename = str(os.path.splitext(os.path.basename(file_path)))

    filename = re.sub(r'^\d+[_\s-]*', '', filename)  # Remove leading numbers
    filename = re.sub(r'[_-]\d{6,8}$', '', filename)  # Remove trailing dates
    filename = re.sub(r'[ \-]+', '_', filename)  # Replace spaces/hyphens
    filename = re.sub(r'[^\w_]', '', filename)  # Remove non-word chars
    return filename.lower()


def process_pdf(file_path):
    doc = pymupdf.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def process_pdf_with_images(file_path, image_output_dir, base_image_name):
    doc = fitz.open(file_path)
    text = ""
    image_count = 0
    os.makedirs(image_output_dir, exist_ok=True)

    for page_num, page in enumerate(doc, start=1):
        text += f"\n--- Page {page_num} ---\n"
        text += page.get_text()

        # Extract images
        for img_index, img in enumerate(page.get_images(full=True), start=1):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_filename = f"{base_image_name}_page{page_num}_img{img_index}.{image_ext}"
            image_filename = image_filename.strip().replace("\n", "").replace("\r", "")
            image_path = image_output_dir / image_filename

            # Save temporarily to check
            with open(image_path, "wb") as f:
                f.write(image_bytes)

            # Check if it's a Verztec logo
            verzy = is_verztec_logo(image_path, threshold=5)
            print(f"[LOG] {image_filename} verztec logo: {verzy}")

            if verzy:
                os.remove(image_path)
                continue  # Skip saving and tagging this image

            # If passed, include image tag
            text += f"<|image_start|>{image_filename}<|image_end|>"

    doc.close()
    return text


def is_verztec_logo(image_path, threshold=5):
    """
    Compare an image against all images in a folder.
    Returns True if any image is similar within the threshold.
    """
    try:
        main_hash = imagehash.phash(Image.open(image_path))
    except Exception as e:
        print(f"Error opening {image_path}: {e}")
        return False

    for filename in os.listdir(vertztec_collection):
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif')):
            compare_path = os.path.join(vertztec_collection, filename)
            try:
                compare_hash = imagehash.phash(Image.open(compare_path))
                diff = main_hash - compare_hash
                if diff <= threshold:
                    return True
            except Exception as e:
                print(f"Error comparing with {compare_path}: {e}")

    return False


def extract_text_and_images_from_doc(file_path, image_output_dir, base_image_name):
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False
    doc = word.Documents.Open(file_path)

    text = ""
    image_count = 0
    os.makedirs(image_output_dir, exist_ok=True)

    for i, inline_shape in enumerate(doc.InlineShapes):
        # Copy the inline image to clipboard
        inline_shape.Select()
        word.Selection.CopyAsPicture()

        try:
            image = ImageGrab.grabclipboard()
            if image:
                image_filename = f"{base_image_name}_img{image_count+1}.png"
                image_path = image_output_dir / image_filename

                # Save temporarily to check logo
                image.save(image_path)
                verzy = is_verztec_logo(image_path, threshold=5)
                print(f"[LOG] {image_filename} verztec logo: {verzy}")

                if verzy:
                    os.remove(image_path)  # Remove unwanted logo image
                    continue  # Skip adding to text or increasing count

                # If passed, include in output
                image_count += 1
                text += f"<|image_start|>{image_filename}<|image_end|>"

        except Exception as e:
            print(f"Failed to extract image {i+1} from Word: {e}")

    # Add the document's text content
    text += doc.Content.Text

    doc.Close(False)
    word.Quit()
    return text

def extract_text_and_images_from_docx(file_path, image_output_dir, base_image_name):
    document = Document(file_path)
    text = ""
    image_count = 0

    # Extract text from paragraphs
    for para in document.paragraphs:
        if para.text.strip():
            text += para.text.strip() + "\n"

    # Extract text from tables (optional, remove if not needed)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    text += cell_text + "\n"

    # Extract images directly from .docx (it's a ZIP archive)
    with zipfile.ZipFile(file_path, 'r') as docx_zip:
        for name in docx_zip.namelist():
            if name.startswith("word/media/"):
                image_ext = os.path.splitext(name)[-1]
                image_count += 1
                image_filename = f"{base_image_name}_img{image_count}{image_ext}"
                image_filename = image_filename.strip().replace("\n", "")
                image_path = image_output_dir / image_filename
                verzy =is_verztec_logo(image_path, threshold=5)
                print(f"[LOG] {image_filename} verztec logo: {verzy}")

                with open(image_path, 'wb') as f:
                    f.write(docx_zip.read(name))

                # Append machine-readable tag
                text += f"<|image_start|>{image_filename}<|image_end|>"

    return text.strip()

def extract_text_and_images_from_pptx(file_path, image_output_dir, base_image_name):
    prs = Presentation(file_path)
    content_lines = []
    image_count = 0
    os.makedirs(image_output_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides, start=1):
        content_lines.append(f"\n--- Slide {slide_num} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                content_lines.append(shape.text.strip())

            if shape.shape_type == 13:  # 13 = PICTURE
                image_count += 1
                image_blob = shape.image.blob
                image_format = shape.image.ext
                image_filename = f"{base_image_name}_slide{slide_num}_img{image_count}.{image_format}"
                image_path = image_output_dir / image_filename

                # Temporarily save image
                with open(image_path, "wb") as img_file:
                    img_file.write(image_blob)

                # Check for Verztec logo
                verzy = is_verztec_logo(image_path, threshold=5)
                print(f"[LOG] {image_filename} verztec logo: {verzy}")

                if verzy:
                    os.remove(image_path)
                    continue  # Skip saving and referencing this image

                # Only include image reference if not Verztec
                content_lines.append(f"<|image_start|>{image_filename}<|image_end|>")

    return "\n\n".join(content_lines)

def extract_inline_content(docx_path, image_dir, base_image_name=None):

    # 1. Prepare image relationships mapping and save images to disk
    image_map = {}  # Map rId -> generated filename
    os.makedirs(image_dir, exist_ok=True)
    image_count = 0

    with zipfile.ZipFile(docx_path) as docx_zip:
        rels_xml = docx_zip.read("word/_rels/document.xml.rels")
        rel_root = ET.fromstring(rels_xml)
        rel_ns = {'r': "http://schemas.openxmlformats.org/package/2006/relationships"}

        for rel in rel_root.findall("r:Relationship", rel_ns):
            r_id = rel.get("Id")
            r_type = rel.get("Type")
            target = rel.get("Target")

            if r_type and r_type.endswith("/image"):
                image_ext = os.path.splitext(target)[-1]
                image_count += 1
                image_filename = f"{base_image_name}_img{image_count}{image_ext}"
                image_path = os.path.join(image_dir, image_filename)

                # Read image data once
                image_data = docx_zip.read(f"word/{target}")

                # Temporarily save to disk to check
                with open(image_path, "wb") as temp_file:
                    temp_file.write(image_data)

                # Check if it's a Verztec logo
                verzy = is_verztec_logo(image_path, threshold=5)
                print(f"[LOG] {image_filename} verztec logo: {verzy}")

                if verzy:
                    os.remove(image_path)  # Delete temp file
                    continue  # Skip this image entirely

                # Image is allowed: keep file and register in image_map
                image_map[r_id] = image_filename

    # 2. Load document for structure
    doc = Document(docx_path)
    body = doc.element.body

    ns = {
        'w':  "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        'wp': "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
        'a':  "http://schemas.openxmlformats.org/drawingml/2006/main",
        'pic':"http://schemas.openxmlformats.org/drawingml/2006/picture",
        'r':  "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        'v':  "urn:schemas-microsoft-com:vml"
    }

    def process_paragraph(paragraph_elem):
        parts = []
        for child in paragraph_elem.iterchildren():
            tag = child.tag.split('}')[-1]
            if tag == "r":
                run = child
            elif tag == "hyperlink":
                for sub in child.iterchildren():
                    if sub.tag.split('}')[-1] == "r":
                        run = sub
                    else:
                        continue
            else:
                continue

            for node in run:
                ntag = node.tag.split('}')[-1]
                if ntag == "t":
                    parts.append(node.text)
                elif ntag == "tab":
                    parts.append("\t")
                elif ntag == "br":
                    parts.append("\n")
                elif ntag == "drawing":
                    blip = node.find(".//a:blip", ns)
                    if blip is not None:
                        embed_id = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                        if embed_id and embed_id in image_map:
                            parts.append(f"<|image_start|>{image_map[embed_id]}<|image_end|>")
                elif ntag == "pict":
                    imdata = node.find(".//v:imagedata", ns)
                    if imdata is not None:
                        embed_id = imdata.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
                        if embed_id and embed_id in image_map:
                            parts.append(f"<|image_start|>{image_map[embed_id]}<|image_end|>")
        return "".join(parts)

    output_lines = []
    for elem in body.iterchildren():
        tag = elem.tag.split('}')[-1]
        if tag == "p":
            output_lines.append(process_paragraph(elem))
        elif tag == "tbl":
            table = docx.table.Table(elem, doc)
            for row in table.rows:
                row_cells_text = []
                for cell in row.cells:
                    cell_text_parts = [process_paragraph(p._element) for p in cell.paragraphs]
                    row_cells_text.append(" ".join(cell_text_parts))
                output_lines.append("\t".join(row_cells_text))

    return "\n".join(output_lines)


# === Main function to process a single file ===
def process_single_file(file_path, images_dir, cleaned_dir, vertztec_collection):
    file_path = Path(file_path)
    ext = file_path.suffix.lower()
    base_filename = file_path.stem
    cleaned_filename = clean_filename_for_embedding(base_filename).strip().replace("\n", "")
    cleaned_text_path = cleaned_dir / (cleaned_filename + ".txt")
    try:
        os.makedirs(images_dir, exist_ok=True)
        os.makedirs(cleaned_dir, exist_ok=True)
        if ext == ".pdf":
            text = process_pdf_with_images(file_path, images_dir, cleaned_filename)
        elif ext == ".docx":
            text = extract_inline_content(file_path, images_dir, base_image_name=cleaned_filename)
        elif ext == ".doc":
            text = extract_text_and_images_from_doc(file_path, images_dir, cleaned_filename)
        elif ext == ".pptx":
            text = extract_text_and_images_from_pptx(file_path, images_dir, cleaned_filename)
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        else:
            raise ValueError(f"Unsupported file type: {ext}")
        # Clean text content
        text = ftfy.fix_text(text)
        text = re.sub(r'\s+', ' ', text).strip()
        text = re.sub(r'\s*(\d+\.)', r'\n\1', text)
        # Save cleaned text
        with open(cleaned_text_path, "w", encoding="utf-8") as out_file:
            out_file.write(text)
        print(f"[INFO] Processed and saved: {cleaned_text_path}")
        return {
            "original_path": str(file_path),
            "cleaned_text_path": str(cleaned_text_path),
            "filename": base_filename,
            "cleaned_filename": cleaned_filename,
            "text_content": text,
            "success": True,
            "error": None,
        }
    except Exception as e:
        print(f"[ERROR] Failed to process {file_path}: {e}")
        return {
            "original_path": str(file_path),
            "cleaned_text_path": None,
            "filename": base_filename,
            "cleaned_filename": cleaned_filename,
            "text_content": None,
            "success": False,
            "error": str(e),
        }

# 
#PDF file cleaning 
# pdf_files = get_all_file_paths(pdf_dir)
# word_files= get_all_file_paths(docx_dir)
# pptx_files= get_all_file_paths(pptx_dir)
# print(f"[LOG] Number of pdf files  : {len(pdf_files)}")
# print(f"[LOG] Number of word files : {len(word_files)}")
# print(f"[LOG] Number of pptx files : {len(pptx_files)}")
# Counters
# success_count = 0
# fail_count = 0
# print('[LOG] starting PDF files')
# for file_path in pdf_files:
   
    
    # Clean filename for embedding
    # cleaned_filename = clean_filename_for_embedding(file_path)
    # base_filename = os.path.splitext(os.path.basename(file_path))[0]

    # Clean filename for both embedding and image naming
    # cleaned_filename = clean_filename_for_embedding(base_filename)
    # cleaned_filename = cleaned_filename.strip().replace("\n", "")

    # This is safe to use for image filenames
    # image_filename = cleaned_filename


    
    # Process the file (either PDF or DOCX)
    # text =process_pdf_with_images(file_path, images_dir, image_filename)
    
    # Apply text cleaning
    # text = ftfy.fix_text(text)  # Fix broken encoding from PDFs or DOCX
    # text = re.sub(r'\s+', ' ', text).strip()  # Remove extra line breaks and normalize spaces
    # text = re.sub(r'\s*(\d+\.)', r'\n\1', text)  # Clean up bullet numbers

    # Output file path for saving cleaned content
    # file_address = cleaned_dir / (cleaned_filename + ".txt")
   

    # Write the extracted text into the cleaned text file
    # with open(file_address, "w", encoding="utf-8") as out:
        # out.write(text)
    # 
    # print(f"[INFO] Finished processing and saved to: {file_address}")


# print('[LOG] starting WORD files')
# for file_path in word_files:
    # mime_type, _ = mimetypes.guess_type(file_path)
    # ext = os.path.splitext(file_path)[-1].lower()
    
    # base_filename = os.path.basename(file_path)
    
    # filename_without_ext = os.path.splitext(base_filename)[0]
    

    # Now fully clean the filename
    # modified_filename = clean_filename_for_embedding(filename_without_ext)
    
    # file_address = cleaned_dir / (modified_filename+ ".txt")
    # base_filename = os.path.splitext(os.path.basename(file_path))[0]

    # Clean filename for both embedding and image naming
    # cleaned_filename = clean_filename_for_embedding(base_filename)
    # cleaned_filename = cleaned_filename.strip().replace("\n", "")

    # This is safe to use for image filenames
    # image_filename = cleaned_filename
    # 
    # print(f"[DEBUG] cleaned_filename = '{cleaned_filename}'")
# 
#    
    # try:
        # if ext == '.docx':
            # 
            # text=extract_inline_content(file_path,images_dir,base_image_name=cleaned_filename)
        #    
        # elif ext == '.doc':
            # text = extract_text_and_images_from_doc(file_path, images_dir, image_filename)
            # text=extract_inline_content(file_path)
            # # 
        # else:
            # print(f"[WARN] Skipped non-Word file: {file_path} (type: {mime_type})")
            # fail_count += 1
            # continue
# 
        #print(f"--- Content from {os.path.basename(file_path)} ---")
        
        # out = open(file_address, "wb") # create a text output
        # 
        # out.write(text.encode("utf8")) # write text of page encoded in utf-8
        # 
        # out.write(bytes((12,))) # write page delimiter (form feed 0x0C)
        # out.close()
    # 
        # 
    #   
        # success_count += 1
        # print(f"[INFO] Finished processing and saved to: {file_address}")
# 
    # except Exception as e:
        # print(f"[WARN] Failed to extract from {file_path}: {e}")
        # fail_count += 1
        # 
# 
# print('[LOG] starting PPTX files')
# for file_path in pptx_files:
    # mime_type, _ = mimetypes.guess_type(file_path)
    # ext = os.path.splitext(file_path)[-1].lower()
    # base_filename = os.path.basename(file_path)
    # filename_without_ext = os.path.splitext(base_filename)[0]
# 
    # modified_filename = clean_filename_for_embedding(filename_without_ext)
    # file_address = cleaned_dir / (modified_filename + ".txt")
# 
    # try:
        # if ext == '.pptx':
            # text = extract_text_and_images_from_pptx(file_path, images_dir, modified_filename)
        # else:
            # print(f"[WARN] Skipped non-PPTX file: {file_path} (type: {mime_type})")
            # fail_count += 1
            # continue
# 
        # with open(file_address, "wb") as out:
            # out.write(text.encode("utf8"))
            # out.write(bytes((12,)))  # Optional delimiter
# 
        # success_count += 1
        # print(f"[INFO] Finished processing and saved to: {file_address}")
# 
    # except Exception as e:
        # print(f"[WARN] Failed to extract from {file_path}: {e}")
        # fail_count += 1
# 
# print(f'[LOG] FAIL COUNT:{fail_count}')
# 


# === Example usage ===
if __name__ == "__main__":
    file_to_process = r"C:\Users\ethan\OneDrive\Desktop\ICP_Verztec_Chatbot-4\data\pdf\3_Offboarding Process on Clean Desk Policy_150125.pdf"
    result = process_single_file(file_to_process, images_dir, cleaned_dir, vertztec_collection)
    if result["success"]:
        print(f"Successfully processed '{result['filename']}'")
        print(f"Cleaned text saved at: {result['cleaned_text_path']}")
    else:
        print(f"Failed to process '{result['filename']}': {result['error']}")
