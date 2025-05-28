import os
import re
import mimetypes
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

import ftfy
import fitz  # pymupdf
import spacy
import win32com.client
from PIL import Image, ImageGrab
import imagehash
from pptx import Presentation
from docx import Document
from docx.oxml.ns import qn

from dotenv import load_dotenv

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import TextLoader, PyPDFLoader
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.schema import Document as langDocument
from langchain_groq import ChatGroq
from langchain_core.output_parsers import StrOutputParser

from MySQLDatabase.Inserting_data import insert_knowledge_chunks


# === Define directories (update as needed) ===
root_dir = Path(os.getcwd())
data_dir = root_dir / "data"
images_dir = data_dir / "images"
pdf_dir = data_dir / "pdf"
docx_dir = data_dir / "word"
pptx_dir = data_dir / "pptx"
cleaned_dir = data_dir / "cleaned"
vertztec_collection = data_dir / "verztec_logo"

# === Helper Functions ===

def get_all_file_paths(directory, extensions=[".pdf", ".doc", ".docx", ".txt", ".pptx"]):
    file_paths = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            if any(file.lower().endswith(ext) for ext in extensions):
                file_paths.append(os.path.abspath(os.path.join(root, file)))
    return file_paths

def clean_filename_for_embedding(file_path):
    filename = str(os.path.splitext(os.path.basename(file_path))[0])
    filename = re.sub(r'^\d+[_\s-]*', '', filename)  # Remove leading numbers
    filename = re.sub(r'[_-]\d{6,8}$', '', filename)  # Remove trailing dates
    filename = re.sub(r'[ \-]+', '_', filename)       # Replace spaces/hyphens
    filename = re.sub(r'[^\w_]', '', filename)        # Remove non-word chars
    return filename.lower()

def is_verztec_logo(image_path, threshold=5):
    try:
        main_hash = imagehash.phash(Image.open(image_path))
    except Exception as e:
        print(f"Error opening {image_path}: {e}")
        return False
    for filename in os.listdir(vertztec_collection):
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif')):
            compare_path = os.path.join(vertztec_collection, filename)
            try:
                compare_hash = imagehash.phash(Image.open(compare_path))
                diff = main_hash - compare_hash
                if diff <= threshold:
                    return True
            except Exception as e:
                print(f"Error comparing with {compare_path}: {e}")
    return False

def process_pdf_with_images(file_path, image_output_dir, base_image_name):
    doc = fitz.open(file_path)
    text = ""
    os.makedirs(image_output_dir, exist_ok=True)
    for page_num, page in enumerate(doc, start=1):
        text += f"\n--- Page {page_num} ---\n"
        text += page.get_text()
        for img_index, img in enumerate(page.get_images(full=True), start=1):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_filename = f"{base_image_name}_page{page_num}_img{img_index}.{image_ext}"
            image_filename = image_filename.strip().replace("\n", "").replace("\r", "")
            image_path = image_output_dir / image_filename
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            verzy = is_verztec_logo(image_path, threshold=5)
            print(f"[LOG] {image_filename} verztec logo: {verzy}")
            if verzy:
                os.remove(image_path)
                continue
            text += f"<|image_start|>{image_filename}<|image_end|>"
    doc.close()
    return text

def extract_text_and_images_from_doc(file_path, image_output_dir, base_image_name):
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False
    doc = word.Documents.Open(str(file_path))
    text = ""
    image_count = 0
    os.makedirs(image_output_dir, exist_ok=True)
    for i, inline_shape in enumerate(doc.InlineShapes):
        inline_shape = doc.InlineShapes.Item(i+1)
        inline_shape.Select()
        word.Selection.CopyAsPicture()
        try:
            image = ImageGrab.grabclipboard()
            if image:
                image_filename = f"{base_image_name}_img{image_count+1}.png"
                image_path = image_output_dir / image_filename
                image.save(image_path)
                verzy = is_verztec_logo(image_path, threshold=5)
                print(f"[LOG] {image_filename} verztec logo: {verzy}")
                if verzy:
                    os.remove(image_path)
                    continue
                image_count += 1
                text += f"<|image_start|>{image_filename}<|image_end|>"
        except Exception as e:
            print(f"Failed to extract image {i+1} from Word: {e}")
    text += doc.Content.Text
    doc.Close(False)
    word.Quit()
    return text

def extract_inline_content(docx_path, image_dir, base_image_name=None):
    image_map = {}
    os.makedirs(image_dir, exist_ok=True)
    image_count = 0
    with zipfile.ZipFile(docx_path) as docx_zip:
        rels_xml = docx_zip.read("word/_rels/document.xml.rels")
        rel_root = ET.fromstring(rels_xml)
        rel_ns = {'r': "http://schemas.openxmlformats.org/package/2006/relationships"}
        for rel in rel_root.findall("r:Relationship", rel_ns):
            r_id = rel.get("Id")
            r_type = rel.get("Type")
            target = rel.get("Target")
            if r_type and r_type.endswith("/image"):
                image_ext = os.path.splitext(target)[-1]
                image_count += 1
                image_filename = f"{base_image_name}_img{image_count}{image_ext}"
                image_path = os.path.join(image_dir, image_filename)
                image_data = docx_zip.read(f"word/{target}")
                with open(image_path, "wb") as temp_file:
                    temp_file.write(image_data)
                verzy = is_verztec_logo(image_path, threshold=5)
                print(f"[LOG] {image_filename} verztec logo: {verzy}")
                if verzy:
                    os.remove(image_path)
                    continue
                image_map[r_id] = image_filename
    doc = Document(docx_path)
    body = doc.element.body
    ns = {
        'w':  "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        'wp': "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
        'a':  "http://schemas.openxmlformats.org/drawingml/2006/main",
        'pic':"http://schemas.openxmlformats.org/drawingml/2006/picture",
        'r':  "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        'v':  "urn:schemas-microsoft-com:vml"
    }
    def process_paragraph(paragraph_elem):
        parts = []
        for child in paragraph_elem.iterchildren():
            tag = child.tag.split('}')[-1]
            if tag == "r":
                run = child
            elif tag == "hyperlink":
                for sub in child.iterchildren():
                    if sub.tag.split('}')[-1] == "r":
                        run = sub
                    else:
                        continue
            else:
                continue
            for node in run:
                ntag = node.tag.split('}')[-1]
                if ntag == "t":
                    parts.append(node.text)
                elif ntag == "tab":
                    parts.append("\t")
                elif ntag == "br":
                    parts.append("\n")
                elif ntag == "drawing":
                    blip = node.find(".//a:blip", ns)
                    if blip is not None:
                        embed_id = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                        if embed_id and embed_id in image_map:
                            parts.append(f"<|image_start|>{image_map[embed_id]}<|image_end|>")
                elif ntag == "pict":
                    imdata = node.find(".//v:imagedata", ns)
                    if imdata is not None:
                        embed_id = imdata.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
                        if embed_id and embed_id in image_map:
                            parts.append(f"<|image_start|>{image_map[embed_id]}<|image_end|>")
        return "".join(parts)
    output_lines = []
    for elem in body.iterchildren():
        tag = elem.tag.split('}')[-1]
        if tag == "p":
            output_lines.append(process_paragraph(elem))
        elif tag == "tbl":
            table = docx.table.Table(elem, doc)
            for row in table.rows:
                row_cells_text = []
                for cell in row.cells:
                    cell_text_parts = [process_paragraph(p._element) for p in cell.paragraphs]
                    row_cells_text.append(" ".join(cell_text_parts))
                output_lines.append("\t".join(row_cells_text))
    return "\n".join(output_lines)

def extract_text_and_images_from_pptx(file_path, image_output_dir, base_image_name):
    prs = Presentation(file_path)
    content_lines = []
    image_count = 0
    os.makedirs(image_output_dir, exist_ok=True)
    for slide_num, slide in enumerate(prs.slides, start=1):
        content_lines.append(f"\n--- Slide {slide_num} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                content_lines.append(shape.text.strip())
            if shape.shape_type == 13:  # PICTURE
                image_count += 1
                image_blob = shape.image.blob
                image_format = shape.image.ext
                image_filename = f"{base_image_name}_slide{slide_num}_img{image_count}.{image_format}"
                image_path = image_output_dir / image_filename
                with open(image_path, "wb") as img_file:
                    img_file.write(image_blob)
                verzy = is_verztec_logo(image_path, threshold=5)
                print(f"[LOG] {image_filename} verztec logo: {verzy}")
                if verzy:
                    os.remove(image_path)
                    continue
                content_lines.append(f"<|image_start|>{image_filename}<|image_end|>")
    return "\n\n".join(content_lines)

# === Main function to process a single file ===
def process_single_file(file_path, images_dir, cleaned_dir, vertztec_collection):
    file_path = Path(file_path)
    ext = file_path.suffix.lower()
    base_filename = file_path.stem
    cleaned_filename = clean_filename_for_embedding(base_filename).strip().replace("\n", "")
    cleaned_text_path = cleaned_dir / (cleaned_filename + ".txt")
    try:
        os.makedirs(images_dir, exist_ok=True)
        os.makedirs(cleaned_dir, exist_ok=True)
        if ext == ".pdf":
            text = process_pdf_with_images(file_path, images_dir, cleaned_filename)
        elif ext == ".docx":
            text = extract_inline_content(file_path, images_dir, base_image_name=cleaned_filename)
        elif ext == ".doc":
            text = extract_text_and_images_from_doc(file_path, images_dir, cleaned_filename)
        elif ext == ".pptx":
            text = extract_text_and_images_from_pptx(file_path, images_dir, cleaned_filename)
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        else:
            raise ValueError(f"Unsupported file type: {ext}")
        # Clean text content
        text = ftfy.fix_text(text)
        text = re.sub(r'\s+', ' ', text).strip()
        text = re.sub(r'\s*(\d+\.)', r'\n\1', text)
        # Save cleaned text
        with open(cleaned_text_path, "w", encoding="utf-8") as out_file:
            out_file.write(text)
        print(f"[INFO] Processed and saved: {cleaned_text_path}")
        return {
            "original_path": file_path,
            "cleaned_text_path": cleaned_text_path,
            "filename": base_filename,
            "cleaned_filename": cleaned_filename,
            "text_content": text,
            "success": True,
            "error": None,
        }
    except Exception as e:
        print(f"[ERROR] Failed to process {file_path}: {e}")
        return {
            "original_path": file_path,
            "cleaned_text_path": None,
            "filename": base_filename,
            "cleaned_filename": cleaned_filename,
            "text_content": None,
            "success": False,
            "error": str(e),
        }


# === Example usage ===
# if __name__ == "__main__":
    file_to_process = r"C:\Users\ethan\OneDrive\Desktop\ICP_Verztec_Chatbot-4\data\pdf\3_Offboarding Process on Clean Desk Policy_150125.pdf"
    result = process_single_file(file_to_process, images_dir, cleaned_dir, vertztec_collection)
    if result["success"]:
        print(f"Successfully processed '{result['filename']}'")
        print(f"Cleaned text saved at: {result['cleaned_text_path']}")
    else:
        print(f"Failed to process '{result['filename']}': {result['error']}")





#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------






# Load environment variables and spaCy model
load_dotenv()
nlp = spacy.load("en_core_web_sm")

# Initialize models
api_key = "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"  # Set in your .env file
model = "deepseek-r1-distill-llama-70b"
deepseek = ChatGroq(api_key=api_key, model_name=model)
deepseek_chain = deepseek | StrOutputParser()

embedding_model = HuggingFaceEmbeddings(
    model_name="BAAI/bge-large-en-v1.5",
    encode_kwargs={'normalize_embeddings': True}
)   

# Helper: Get all file paths recursively

def get_all_file_paths(directory, extensions=[".pdf", ".doc", ".docx", ".txt"]):
    file_paths = []
    for root, _, files in os.walk(directory):
        for file in files:
            if any(file.lower().endswith(ext) for ext in extensions):
                file_paths.append(os.path.abspath(os.path.join(root, file)))
    return file_paths

# Text splitters
fallback_splitter = RecursiveCharacterTextSplitter(chunk_size=250, chunk_overlap=50)
bullet_pattern = r'(?m)(?:^\s*(\d+[\.\)]|[a-zA-Z][\.\)]|[-\u2022\*])\s+)(.*?)(?=^\s*(\d+[\.\)]|[a-zA-Z][\.\)]|[-\u2022\*])\s+|$)'

def spacy_fallback_splitter(text, chunk_size=400, overlap=50):
    doc = nlp(text)
    sentences = [sent.text.strip() for sent in doc.sents if sent.text.strip()]
    chunks, current_chunk = [], ""
    for sentence in sentences:
        if len(current_chunk) + len(sentence) <= chunk_size:
            current_chunk += " " + sentence
        else:
            chunks.append(current_chunk.strip())
            current_chunk = sentence
    if current_chunk:
        chunks.append(current_chunk.strip())

    if overlap > 0:
        overlapped_chunks = []
        for i in range(len(chunks)):
            start = max(0, i - 1)
            merged = " ".join(chunks[start:i+1])
            overlapped_chunks.append(merged.strip())
        return overlapped_chunks

    return chunks

def split_text(text):
    bullet_pattern = r'(?m)(^\s*(\d+[\.\)]|[a-zA-Z][\.\)]|[-\u2022\*])\s+.*(?:\n(?!^\s*(\d+[\.\)]|[a-zA-Z][\.\)]|[-\u2022\*])\s+).*)*)'
    matches = list(re.finditer(bullet_pattern, text, flags=re.MULTILINE))
    chunks = [m.group(0).strip() for m in matches if len(m.group(0).strip()) > 5 and not m.group(0).strip().isdigit()]
    
    final_chunks = []
    for chunk in chunks:
        final_chunks.extend(fallback_splitter.split_text(chunk) if len(chunk) > 600 else [chunk])
    
    return final_chunks or spacy_fallback_splitter(text)


def replace_image_tags_with_placeholders(text):
    pattern = r"<\|image_start\|>(.*?)<\|image_end\|>"
    matches = re.findall(pattern, text)
    replacements = {}
    image_name_map = {}
    for i, match in enumerate(matches):
        key = f"[[IMAGE_TAG_{i}]]"
        full_tag = f"<|image_start|>{match}<|image_end|>"
        text = text.replace(full_tag, key)
        replacements[key] = full_tag
        image_name_map[key] = match  # Just the image file name
    return text, replacements, image_name_map


def restore_image_placeholders_and_collect_metadata(chunks, replacements, image_name_map):
    restored_chunks = []
    chunk_image_metadata = []

    for chunk in chunks:
        images_in_chunk = []
        for placeholder, original in replacements.items():
            if placeholder in chunk:
                chunk = chunk.replace(placeholder, original)
                images_in_chunk.append(image_name_map[placeholder])
        restored_chunks.append(chunk)
        chunk_image_metadata.append(images_in_chunk)
    
    return restored_chunks, chunk_image_metadata

# Paths
root_dir = Path.cwd()
data_dir = root_dir / "data/cleaned"
cleaned_files = get_all_file_paths(data_dir)

# Chunk and embed all docs
all_chunks = []
for file_path in cleaned_files:
    indv_chunks = []
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    
    with open(file_path, 'r', encoding='utf-8') as f:
        text = f.read().lower()

    # Generate semantic description with LLM
    query = (
    "Summarize this internal document from Verztec in 50 words. "
    "Describe its purpose, main contents, and the types of internal questions it can help answer. your response should help an embedding model retreive the relevant chunks "
    "Respond concisely and in a single paragraph:\n\n" + text
    )

    raw_response = deepseek_chain.invoke(query)
    cleaned_response = re.sub(r"<think>.*?</think>", "", raw_response, flags=re.DOTALL).strip()
    #cleaned_response='' 
    description = cleaned_response
    print('------------------------------------------')
    print(description)

    # Replace image tags with placeholders and keep track of their names
    text, image_replacements, image_name_map = replace_image_tags_with_placeholders(text)

    # Split text into chunks
    smart_chunks = split_text(text)

    # Restore image tags and get image names per chunk
    smart_chunks, chunk_image_lists = restore_image_placeholders_and_collect_metadata(
        smart_chunks, image_replacements, image_name_map
    )

    # Build langDocument objects
    for i, chunk in enumerate(smart_chunks):
        image_list = chunk_image_lists[i]

        # Enriched chunk for global index
        enriched_chunk = f"[Description: {description}] [Document: {base_name}] {chunk}"
        all_chunks.append(langDocument(
            page_content=enriched_chunk,
            metadata={
                "source": base_name,
                "chunk_id": f"{base_name}_{i}",
                "clean_chunk": chunk,
                "images": image_list
            }
        ))

        # Raw chunk for individual index
        indv_chunks.append(langDocument(
            page_content=chunk,
            metadata={
                "source": base_name,
                "chunk_id": f"{base_name}_{i}",
                "images": image_list
            }
        ))

    # Save FAISS index for individual file
    #indv_db = FAISS.from_documents(indv_chunks, embedding_model)
    #indv_db.save_local(f"indv_doc_faiss/{base_name}_index")

for chunk in all_chunks:
    print (chunk.page_content)
    print('--------------------------------------------------------------------------------------------------')
# Save combined vector index

global_db = FAISS.from_documents(all_chunks, embedding_model)
global_db.save_local("faiss_index3")
# saving chunks in sql 
insert_knowledge_chunks(all_chunks)

from langchain.schema import Document as langDocument
from langchain_community.vectorstores import FAISS
import re

def chunk_and_index_text(text: str, base_name: str, faiss_index_path: str = None):
    """
    Given extracted text and a base_name (filename without extension),
    chunk the text, embed, and add to existing FAISS index (if given) or create new.
    Save updated FAISS index locally.

    Returns updated FAISS index object.
    """

    # 1. Lowercase for consistency (like your original)
    text = text.lower()

    # 2. Generate semantic description with LLM (your Deepseek chain)
    query = (
        "Summarize this internal document from Verztec in 50 words. "
        "Describe its purpose, main contents, and the types of internal questions it can help answer. "
        "Respond concisely and in a single paragraph:\n\n" + text
    )
    raw_response = deepseek_chain.invoke(query)
    cleaned_response = re.sub(r"<think>.*?</think>", "", raw_response, flags=re.DOTALL).strip()
    description = cleaned_response

    # 3. Replace image tags with placeholders
    text, image_replacements, image_name_map = replace_image_tags_with_placeholders(text)

    # 4. Split text into chunks
    smart_chunks = split_text(text)

    # 5. Restore image tags and get images per chunk
    smart_chunks, chunk_image_lists = restore_image_placeholders_and_collect_metadata(
        smart_chunks, image_replacements, image_name_map
    )

    # 6. Build langDocument objects for chunks
    all_chunks = []
    for i, chunk in enumerate(smart_chunks):
        image_list = chunk_image_lists[i]

        enriched_chunk = f"[Description: {description}] [Document: {base_name}] {chunk}"
        all_chunks.append(langDocument(
            page_content=enriched_chunk,
            metadata={
                "source": base_name,
                "chunk_id": f"{base_name}_{i}",
                "clean_chunk": chunk,
                "images": image_list
            }
        ))

    # 7. Load existing FAISS index or create new
    if faiss_index_path and Path(faiss_index_path).exists():
        print(f"Loading existing FAISS index from {faiss_index_path}")
        global_db = FAISS.load_local(faiss_index_path, embedding_model)
        # Add new docs to existing index
        global_db.add_documents(all_chunks)
    else:
        print("Creating new FAISS index")
        global_db = FAISS.from_documents(all_chunks, embedding_model)

    # 8. Save updated index
    save_path = faiss_index_path if faiss_index_path else "faiss_index3"
    global_db.save_local(save_path)
    print(f"Saved FAISS index to {save_path}")

    # 9. Insert chunks into SQL DB (optional, as before)
    insert_knowledge_chunks(all_chunks)

    return global_db


file_path = r"C:\Users\ethan\OneDrive\Desktop\ICP_Verztec_Chatbot-4\data\pdf\3_Offboarding Process on Clean Desk Policy_150125.pdf"
result = process_single_file(file_path)

if result["success"]:
    text = result["text"]
    base_name = result["cleaned_filename"]
    # Call your chunk and indexing function here
    updated_index = chunk_and_index_text(text, base_name, faiss_index_path="faiss_index2")
else:
    print("Extraction failed:", result["error"])
